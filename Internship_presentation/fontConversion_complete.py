from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os

def change_font_and_justify(pptx_file, output_folder):
    prs = Presentation(pptx_file)

    # Define new font details
    new_font_name = 'Times New Roman'
    title_font_size = Pt(38)
    text_font_size = Pt(30)
    new_font_color = RGBColor(0, 0, 0)  # Black color

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            if shape.has_title:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.JUSTIFY
                    for run in paragraph.runs:
                        run.font.name = new_font_name
                        run.font.size = title_font_size
                        run.font.color.rgb = new_font_color
            else:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.JUSTIFY
                    for run in paragraph.runs:
                        run.font.name = new_font_name
                        run.font.size = text_font_size
                        run.font.color.rgb = new_font_color

    # Construct the output file path
    output_file = os.path.join(output_folder, 'FontAndJustifyChanged_presentation.pptx')

    prs.save(output_file)
    print(f"Font and text justified successfully! Saved to: {output_file}")

if __name__ == "__main__":
    pptx_file = r'C:\Users\kamalesh.kb\INTERNSHIP_REPORT\INTERNSHIP_REPORT\internship_presentation.pptx'  # Replace with your actual PPTX file path
    output_folder = r'C:\Users\kamalesh.kb\INTERNSHIP_REPORT\INTERNSHIP_REPORT'  # Replace with your desired output folder

    change_font_and_justify(pptx_file, output_folder)
