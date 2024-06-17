from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os

def justify_text(pptx_file, output_folder):
    prs = Presentation(pptx_file)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            for paragraph in shape.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.JUSTIFY

    # Construct the output file path
    output_file = os.path.join(output_folder, 'justified_presentation.pptx')

    prs.save(output_file)
    print(f"Text justified successfully! Saved to: {output_file}")

if __name__ == "__main__":
    pptx_file = r'C:\Users\kamalesh.kb\INTERNSHIP_REPORT\INTERNSHIP_REPORT\FontChanged_presentation.pptx'  # Replace with your actual PPTX file path
    output_folder = r'C:\Users\kamalesh.kb\INTERNSHIP_REPORT\INTERNSHIP_REPORT'  # Replace with your desired output folder

    justify_text(pptx_file, output_folder)
